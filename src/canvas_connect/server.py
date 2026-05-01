"""Canvas Connect MCP Server - Main server implementation."""

import io
import os
import logging
import re
from typing import Any
from datetime import datetime
from pathlib import Path

from dotenv import load_dotenv
import anthropic
from mcp.server import Server

# Load environment variables from .env file
env_path = Path(__file__).parent.parent.parent / ".env"
load_dotenv(env_path)

# Rubrics for assignments
RUBRICS = {
    1027382: {  # Introductions: Are You Google or Microsoft? Android or iPhone?
        "assignment_name": "Introductions: Are You Google or Microsoft? Android or iPhone?",
        "total_points": 100,
        "dimensions": [
            {
                "name": "Quality",
                "max_points": 25,
                "level_exemplars": {
                    "25": ["Great - Initial post is appropriate and meets word count of 350 words. Post is thoughtful and reflective"],
                    "15": ["Good - Initial post word count less than 350 words."],
                    "10": ["Fair - The post is perfunctory, does not contain any personal experience, leaves little room for more conversation."],
                    "0": ["Unacceptable - No posting"]
                }
            },
            {
                "name": "Engagement",
                "max_points": 25,
                "level_exemplars": {
                    "25": ["Great - Aware of the needs to interact with fellow students; attempts to motivate the group discussion; presents creative approaches to the topics. Interacts freely. Responds to more than the minimum number of colleagues."],
                    "15": ["Good - Presents relevant viewpoints for consideration by group; responds to at least 1 colleague."],
                    "10": ["Fair - Responses to 1 or 2 colleagues is perfunctory and does not add value to the conversation."],
                    "0": ["Poor - No feedback provided to fellow students."]
                }
            },
            {
                "name": "Relevance",
                "max_points": 25,
                "level_exemplars": {
                    "25": ["Great - Initial post is in response to the provided topics and prompts further discussion or inquiry."],
                    "15": ["Good - The initial post just addressed the topics provided with close-ended responses."],
                    "10": ["Fair - Not all of the topics provided were addressed or some of the responses were not relevant."],
                    "0": ["Unacceptable - There was no initial posting or none of the responses were relevant to the topics provided."]
                }
            },
            {
                "name": "Mechanics",
                "max_points": 25,
                "level_exemplars": {
                    "25": ["Great - Clear and articulate; grammar, punctuation, and spelling are correct."],
                    "15": ["Good - Contains minor errors that don't cloud meaning and very few (if any) mechanical errors."],
                    "10": ["Fair - Contains several proofing errors or uses text-messaging shortcuts, making the post hard to follow."],
                    "0": ["Poor - Contains multiple mechanical errors (spelling, grammar, and punctuation); the post is sloppy and uses text-messaging shortcuts; meaning of the post is hard to understand, or there is no initial post."]
                }
            }
        ]
    },
    1027378: {  # DISCUSS: The IoT Discussion
        "assignment_name": "The IoT Discussion",
        "total_points": 100,
        "dimensions": [
            {
                "name": "Initial Post",
                "max_points": 20,
                "level_exemplars": {
                    "20": ["Great - Initial post is appropriate, meets word count of 300 words, and is free of errors."],
                    "12": ["Good - Initial post is appropriate, meets the word count, but multiple errors are evident."],
                    "5": ["Fair - The initial post does not meet the word count of 300 words and/or errors are too prevalent to read content."],
                    "0": ["Unacceptable - No posting."]
                }
            },
            {
                "name": "Response",
                "max_points": 24,
                "level_exemplars": {
                    "24": ["Great - Answered all questions in initial post."],
                    "18": ["Good - Answered 3 questions in the initial post."],
                    "12": ["Fair - Answered two of the questions in the initial post."],
                    "6": ["Fair - Answered 1 question in the initial post."],
                    "0": ["Poor - Has an initial post, but none of the questions are addressed."]
                }
            },
            {
                "name": "Thoroughness of Response",
                "max_points": 32,
                "level_exemplars": {
                    "32": ["Great - All 4 questions were answered, and the responses were very well thought out."],
                    "24": ["Good - Only 3 questions were answered, but the responses were very well thought out OR all questions were answered, but only 3 had thoughtful answers."],
                    "16": ["Fair - Only 2 questions were answered, but the responses were very well thought out OR all questions were answered, but only 2 had thoughtful answers."],
                    "8": ["Fair - Only 1 question was answered, but the response was very well thought out OR all questions were answered, but only 1 had a thoughtful answer."],
                    "0": ["Poor - No questions were addressed OR all questions had substandard answers."]
                }
            },
            {
                "name": "Replies",
                "max_points": 24,
                "level_exemplars": {
                    "24": ["Great - Replied to at least 2 colleagues"],
                    "12": ["Good - Replied to only 1 colleague"],
                    "0": ["Poor - Didn't reply to anyone."]
                }
            }
        ]
    },
    1027381: {  # AI in our Future: A Student's Beginning Guide to Artificial Intelligence
        "assignment_name": "AI in our Future: A Student's Beginning Guide to Artificial Intelligence",
        "total_points": 100,
        "dimensions": [
            {
                "name": "Quality",
                "max_points": 25,
                "level_exemplars": {
                    "25": ["Great - Initial post is appropriate and meets word count of 350 words. Post is thoughtful and reflective"],
                    "15": ["Good - Initial post word count less than 350 words."],
                    "10": ["Fair - The post is perfunctory, does not contain any personal experience, leaves little room for more conversation."],
                    "0": ["Unacceptable - No posting."]
                }
            },
            {
                "name": "Relevance",
                "max_points": 25,
                "level_exemplars": {
                    "25": ["Great - Initial post is in response to the provided topics and prompts further discussion or inquiry."],
                    "15": ["Good - The initial post just addressed the topics provided with close-ended responses."],
                    "10": ["Fair - Not all of the topics provided were addressed or some of the responses were not relevant."],
                    "0": ["Poor - There was no initial posting or none of the responses were relevant to the topics provided."]
                }
            },
            {
                "name": "Engagement",
                "max_points": 25,
                "level_exemplars": {
                    "25": ["Great - Aware of the needs to interact with fellow students; attempts to motivate the group discussion; presents creative approaches to the topics. Interacts freely. Responds to more than the minimum number of colleagues."],
                    "15": ["Good - Presents relevant viewpoints for consideration by group; responds to at least 1 colleague."],
                    "10": ["Fair - Responses to 1 or 2 colleagues is perfunctory and does not add value to the conversation."],
                    "0": ["Unacceptable - There was no initial posting or none of the responses were relevant to the topics provided."]
                }
            },
            {
                "name": "Mechanics",
                "max_points": 25,
                "level_exemplars": {
                    "25": ["Great - Clear and articulate; grammar, punctuation, and spelling are correct."],
                    "15": ["Good - Contains minor errors that don't cloud meaning and very few (if any) mechanical errors."],
                    "10": ["Fair - Contains several proofing errors or uses text-messaging shortcuts, making the post hard to follow."],
                    "0": ["Poor - Contains multiple mechanical errors (spelling, grammar, and punctuation); the post is sloppy and uses text-messaging shortcuts; meaning of the post is hard to understand, or there is no initial post."]
                }
            }
        ]
    }
}
from mcp.types import (
    Tool,
    TextContent,
    ImageContent,
    EmbeddedResource,
)
from canvasapi import Canvas
from canvasapi.course import Course
from canvasapi.exceptions import CanvasException

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("canvas-connect")

# Initialize MCP server
app = Server("canvas-connect")

# Global Canvas instance
canvas_instance: Canvas | None = None
course_instance: Course | None = None


def _strip_html(text: str) -> str:
    """Remove HTML tags from text."""
    if not text:
        return ""
    return re.sub(r'<[^>]+>', '', text).strip()


async def _download_canvas_file(url: str) -> bytes:
    """Download a file from Canvas with API token authentication."""
    import httpx
    api_token = os.getenv("CANVAS_API_TOKEN")
    async with httpx.AsyncClient(follow_redirects=True, timeout=60.0) as client:
        response = await client.get(url, headers={"Authorization": f"Bearer {api_token}"})
        response.raise_for_status()
        return response.content


def _extract_docx_content(file_bytes: bytes) -> str:
    """Extract text and structure from a Word (.docx) document."""
    from docx import Document
    doc = Document(io.BytesIO(file_bytes))

    word_count = sum(len(p.text.split()) for p in doc.paragraphs if p.text.strip())
    lines = [
        f"WORD COUNT: ~{word_count}",
        f"PARAGRAPHS: {len([p for p in doc.paragraphs if p.text.strip()])}",
        f"SECTIONS: {len(doc.sections)}",
        f"TABLES: {len(doc.tables)}",
        "",
        "=== DOCUMENT CONTENT ===",
    ]

    for para in doc.paragraphs:
        if not para.text.strip():
            continue
        style_name = para.style.name if para.style else ""
        if style_name.startswith("Heading"):
            try:
                level = int(style_name.split()[-1])
                lines.append(f"\n{'#' * level} {para.text}")
            except (ValueError, IndexError):
                lines.append(f"\n## {para.text}")
        else:
            lines.append(para.text)

    if doc.tables:
        lines.append(f"\n=== TABLES ({len(doc.tables)} total) ===")
        for i, table in enumerate(doc.tables[:5]):
            lines.append(f"\nTable {i + 1}:")
            for row in table.rows:
                lines.append(" | ".join(cell.text.strip() for cell in row.cells))

    return "\n".join(lines)


def _extract_xlsx_content(file_bytes: bytes) -> str:
    """Extract data and structure from an Excel (.xlsx) workbook."""
    import openpyxl

    wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
    lines = [f"SHEETS: {', '.join(wb.sheetnames)}", ""]

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        lines.append(f"=== SHEET: {sheet_name} ===")
        if ws.dimensions and ws.dimensions != "A1:A1":
            lines.append(f"Data Range: {ws.dimensions}")
        if hasattr(ws, "_charts") and ws._charts:
            lines.append(f"Charts: {len(ws._charts)}")

        row_count = 0
        for row in ws.iter_rows(values_only=True):
            if any(cell is not None for cell in row):
                lines.append(" | ".join(str(v) if v is not None else "" for v in row))
                row_count += 1
                if row_count >= 75:
                    lines.append("... (additional rows truncated)")
                    break
        lines.append("")

    # Second pass to capture formulas (load_workbook with data_only=False)
    try:
        wb_f = openpyxl.load_workbook(io.BytesIO(file_bytes))
        formulas = []
        for sheet_name in wb_f.sheetnames:
            ws = wb_f[sheet_name]
            for row in ws.iter_rows():
                for cell in row:
                    if isinstance(cell.value, str) and cell.value.startswith("="):
                        formulas.append(f"  {sheet_name}!{cell.coordinate}: {cell.value}")
                        if len(formulas) >= 30:
                            break
                if len(formulas) >= 30:
                    break
        if formulas:
            lines.append("=== FORMULAS USED (sample) ===")
            lines.extend(formulas)
    except Exception:
        pass

    return "\n".join(lines)


def _extract_pptx_content(file_bytes: bytes) -> str:
    """Extract text and structure from a PowerPoint (.pptx) presentation."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    lines = [
        f"TOTAL SLIDES: {len(prs.slides)}",
        f"SLIDE DIMENSIONS: {prs.slide_width.inches:.1f}\" x {prs.slide_height.inches:.1f}\"",
        "",
    ]

    for i, slide in enumerate(prs.slides):
        lines.append(f"=== SLIDE {i + 1} ===")
        if slide.slide_layout and slide.slide_layout.name:
            lines.append(f"Layout: {slide.slide_layout.name}")

        title_text = slide.shapes.title.text if slide.shapes.title else None
        if title_text:
            lines.append(f"Title: {title_text}")

        for shape in slide.shapes:
            if not hasattr(shape, "text_frame") or not shape.text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text and text != title_text:
                    lines.append(f"  • {text}")

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                lines.append(f"  [Notes: {notes[:300]}{'...' if len(notes) > 300 else ''}]")

        image_count = sum(1 for shape in slide.shapes if shape.shape_type == 13)
        if image_count:
            lines.append(f"  [Images: {image_count}]")

        lines.append("")

    return "\n".join(lines)


def _get_canvas_rubric(assignment) -> str | None:
    """Extract a rubric from a Canvas assignment object, if one is attached."""
    rubric = getattr(assignment, "rubric", None)
    if not rubric:
        return None

    lines = []
    total_points = 0
    for criterion in rubric:
        pts = criterion.get("points", 0)
        total_points += pts
        lines.append(f"\n{criterion.get('description', 'Criterion')} ({pts} pts)")
        lines.append("-" * 40)
        for rating in sorted(criterion.get("ratings", []), key=lambda r: r.get("points", 0), reverse=True):
            desc = rating.get("description", "")
            long_desc = rating.get("long_description", "")
            label = f"{desc} - {long_desc}" if long_desc else desc
            lines.append(f"  {rating.get('points', 0)} pts: {label}")

    lines.append(f"\nTotal Points: {total_points}")
    return "\n".join(lines)


async def _grade_file_with_claude(
    student_name: str,
    file_name: str,
    file_content: str,
    assignment_name: str,
    assignment_description: str,
    rubric_text: str,
    points_possible: float,
) -> str:
    """Send extracted file content to Claude for rubric-based grading."""
    api_key = os.getenv("ANTHROPIC_API_KEY")
    if not api_key:
        raise ValueError("ANTHROPIC_API_KEY environment variable not set")

    client = anthropic.Anthropic(api_key=api_key)

    max_chars = 8000
    if len(file_content) > max_chars:
        file_content = file_content[:max_chars] + "\n... [content truncated due to length]"

    prompt = f"""You are grading a student's final project file submission for a college-level introductory computer applications course.

ASSIGNMENT: {assignment_name}
STUDENT: {student_name}
FILE: {file_name}
TOTAL POINTS POSSIBLE: {points_possible}

ASSIGNMENT INSTRUCTIONS:
{assignment_description or "(No description provided)"}

GRADING RUBRIC:
{rubric_text or "(No rubric found — grade based on the assignment instructions and general quality)"}

STUDENT'S SUBMISSION CONTENT:
{file_content}

Grade this submission carefully using the rubric. Assign a score for each dimension and explain why. Then give an overall score and constructive feedback.

Format your response EXACTLY as follows:

DIMENSION SCORES:
- [Dimension Name]: [Points Earned] / [Max Points] — [Brief justification]
(repeat for each dimension)

TOTAL SCORE: [X] / {points_possible}
PERCENTAGE: [X]%
LETTER GRADE: [A/B/C/D/F]

STUDENT FEEDBACK:
[3-5 sentences of specific, constructive feedback addressing strengths and areas for improvement.]"""

    message = client.messages.create(
        model="claude-sonnet-4-20250514",
        max_tokens=2048,
        messages=[{"role": "user", "content": prompt}],
    )
    return message.content[0].text


def _rubric_text_from_dict(assignment_id: int) -> str | None:
    """Return a formatted rubric string from the RUBRICS dict, if present."""
    if assignment_id not in RUBRICS:
        return None
    rubric = RUBRICS[assignment_id]
    lines = [f"RUBRIC: {rubric['assignment_name']}", f"Total Points: {rubric['total_points']}"]
    for dim in rubric["dimensions"]:
        lines.append(f"\n{dim['name']} ({dim['max_points']} pts)")
        lines.append("-" * 40)
        for pts in sorted(dim["level_exemplars"].keys(), key=int, reverse=True):
            lines.append(f"  {pts} pts: {dim['level_exemplars'][pts][0]}")
    return "\n".join(lines)


async def _grade_submission(
    course,
    assignment,
    submission,
    student_name: str,
    rubric_text: str | None,
    assignment_description: str,
    points_possible: float,
    auto_update: bool,
) -> list[str]:
    """Download, extract, and grade all file attachments in a single submission."""
    attachments = getattr(submission, "attachments", []) or []
    if not attachments:
        return [f"  No file attachments found — student may not have submitted a file."]

    lines = []
    for attachment in attachments:
        file_name = getattr(attachment, "filename", "unknown")
        file_url = getattr(attachment, "url", None)
        file_size = getattr(attachment, "size", 0)

        lines.append(f"\n  File: {file_name} ({file_size:,} bytes)")

        if not file_url:
            lines.append("  ERROR: No download URL available for this attachment.")
            continue

        try:
            file_bytes = await _download_canvas_file(file_url)
        except Exception as e:
            lines.append(f"  ERROR downloading file: {e}")
            continue

        file_lower = file_name.lower()
        try:
            if file_lower.endswith(".docx"):
                file_content = _extract_docx_content(file_bytes)
                file_type = "Word Document"
            elif file_lower.endswith(".xlsx"):
                file_content = _extract_xlsx_content(file_bytes)
                file_type = "Excel Spreadsheet"
            elif file_lower.endswith(".pptx"):
                file_content = _extract_pptx_content(file_bytes)
                file_type = "PowerPoint Presentation"
            elif file_lower.endswith((".doc", ".xls", ".ppt")):
                lines.append(
                    f"  WARNING: Legacy Office format (.doc/.xls/.ppt) not supported. "
                    "Ask the student to re-submit in the modern format (.docx/.xlsx/.pptx)."
                )
                continue
            else:
                lines.append(f"  WARNING: Unsupported file type — cannot grade '{file_name}'.")
                continue
        except Exception as e:
            lines.append(f"  ERROR extracting content from {file_name}: {e}")
            continue

        lines.append(f"  Type: {file_type}")

        try:
            grade_result = await _grade_file_with_claude(
                student_name=student_name,
                file_name=file_name,
                file_content=file_content,
                assignment_name=assignment.name,
                assignment_description=assignment_description,
                rubric_text=rubric_text or "",
                points_possible=points_possible,
            )
        except Exception as e:
            lines.append(f"  ERROR grading with Claude: {e}")
            continue

        lines.append("")
        lines.append(grade_result)

        if auto_update:
            score_match = re.search(r"TOTAL SCORE:\s*(\d+(?:\.\d+)?)\s*/\s*\d+", grade_result)
            if score_match:
                score = float(score_match.group(1))
                update_data: dict = {"posted_grade": str(score)}
                feedback_match = re.search(
                    r"STUDENT FEEDBACK:\n(.+?)(?:\n\n|\Z)", grade_result, re.DOTALL
                )
                if feedback_match:
                    update_data["comment"] = {"text_comment": feedback_match.group(1).strip()}
                try:
                    submission.edit(submission=update_data)
                    lines.append(f"\n  Grade updated in Canvas: {score}/{points_possible}")
                except Exception as e:
                    lines.append(f"\n  Could not update grade in Canvas: {e}")
            else:
                lines.append("\n  Could not parse score from Claude response to auto-update grade.")

    return lines


def get_canvas() -> Canvas:
    """Get or create Canvas API instance."""
    global canvas_instance
    if canvas_instance is None:
        api_url = os.getenv("CANVAS_API_URL")
        api_token = os.getenv("CANVAS_API_TOKEN")

        if not api_url or not api_token:
            raise ValueError(
                "Missing Canvas API credentials. Please set CANVAS_API_URL and CANVAS_API_TOKEN environment variables."
            )

        canvas_instance = Canvas(api_url, api_token)
        logger.info(f"Connected to Canvas at {api_url}")

    return canvas_instance


def get_course() -> Course:
    """Get or create Course instance."""
    global course_instance
    if course_instance is None:
        canvas = get_canvas()
        course_id = os.getenv("CANVAS_COURSE_ID")

        if not course_id:
            raise ValueError(
                "Missing course ID. Please set CANVAS_COURSE_ID environment variable."
            )

        course_instance = canvas.get_course(int(course_id))
        logger.info(f"Loaded course: {course_instance.name}")

    return course_instance


@app.list_tools()
async def list_tools() -> list[Tool]:
    """List all available Canvas management tools."""
    return [
        # Assignment tools
        Tool(
            name="list_assignments",
            description="List all assignments in the course with their details (name, due date, points, published status)",
            inputSchema={
                "type": "object",
                "properties": {},
            },
        ),
        Tool(
            name="create_assignment",
            description="Create a new assignment in the course",
            inputSchema={
                "type": "object",
                "properties": {
                    "name": {"type": "string", "description": "Assignment name"},
                    "description": {"type": "string", "description": "Assignment description (HTML supported)"},
                    "points": {"type": "number", "description": "Points possible"},
                    "due_at": {"type": "string", "description": "Due date in ISO format (e.g., 2024-12-31T23:59:59)"},
                    "published": {"type": "boolean", "description": "Whether to publish immediately", "default": False},
                },
                "required": ["name", "points"],
            },
        ),
        Tool(
            name="get_assignment_submissions",
            description="Get all submissions for a specific assignment",
            inputSchema={
                "type": "object",
                "properties": {
                    "assignment_id": {"type": "number", "description": "Assignment ID"},
                },
                "required": ["assignment_id"],
            },
        ),
        Tool(
            name="update_grade",
            description="Update a student's grade for an assignment",
            inputSchema={
                "type": "object",
                "properties": {
                    "assignment_id": {"type": "number", "description": "Assignment ID"},
                    "user_id": {"type": "number", "description": "Student user ID"},
                    "grade": {"type": "string", "description": "Grade (number or letter)"},
                    "comment": {"type": "string", "description": "Optional comment"},
                },
                "required": ["assignment_id", "user_id", "grade"],
            },
        ),

        # Student tools
        Tool(
            name="list_students",
            description="List all students enrolled in the course",
            inputSchema={
                "type": "object",
                "properties": {},
            },
        ),
        Tool(
            name="get_student_grades",
            description="Get a student's grades across all assignments",
            inputSchema={
                "type": "object",
                "properties": {
                    "user_id": {"type": "number", "description": "Student user ID"},
                },
                "required": ["user_id"],
            },
        ),
        Tool(
            name="get_grade_summary",
            description="Get a summary of all students with their current total grade/score in the course",
            inputSchema={
                "type": "object",
                "properties": {},
            },
        ),

        # Announcement tools
        Tool(
            name="list_announcements",
            description="List recent announcements in the course",
            inputSchema={
                "type": "object",
                "properties": {
                    "limit": {"type": "number", "description": "Number of announcements to retrieve", "default": 10},
                },
            },
        ),
        Tool(
            name="create_announcement",
            description="Create a new announcement in the course",
            inputSchema={
                "type": "object",
                "properties": {
                    "title": {"type": "string", "description": "Announcement title"},
                    "message": {"type": "string", "description": "Announcement message (HTML supported)"},
                    "published": {"type": "boolean", "description": "Whether to publish immediately", "default": True},
                },
                "required": ["title", "message"],
            },
        ),

        # Module tools
        Tool(
            name="list_modules",
            description="List all modules in the course",
            inputSchema={
                "type": "object",
                "properties": {},
            },
        ),
        Tool(
            name="get_module_items",
            description="Get all items in a specific module",
            inputSchema={
                "type": "object",
                "properties": {
                    "module_id": {"type": "number", "description": "Module ID"},
                },
                "required": ["module_id"],
            },
        ),
        Tool(
            name="create_module",
            description="Create a new module in the course",
            inputSchema={
                "type": "object",
                "properties": {
                    "name": {"type": "string", "description": "Module name"},
                    "published": {"type": "boolean", "description": "Whether to publish immediately", "default": False},
                },
                "required": ["name"],
            },
        ),

        # Rubric tools
        Tool(
            name="get_rubric",
            description="Get the grading rubric for a specific assignment. Returns the rubric criteria, point values, and level descriptions to use when grading.",
            inputSchema={
                "type": "object",
                "properties": {
                    "assignment_id": {"type": "number", "description": "Assignment ID to get rubric for"},
                },
                "required": ["assignment_id"],
            },
        ),

        # Discussion tools
        Tool(
            name="get_discussion_posts",
            description="Get all posts and replies for a discussion topic. Returns the initial posts from each student and their replies to other students, with word counts.",
            inputSchema={
                "type": "object",
                "properties": {
                    "topic_id": {"type": "number", "description": "Discussion assignment ID or discussion topic ID (for graded discussions, use the assignment ID)"},
                    "user_id": {"type": "number", "description": "Optional: Filter to show only posts from a specific student"},
                },
                "required": ["topic_id"],
            },
        ),

        # AI Detection tools
        Tool(
            name="check_ai_writing",
            description="Analyze a student's discussion post to detect if it was likely written by AI. Returns a likelihood assessment and indicators.",
            inputSchema={
                "type": "object",
                "properties": {
                    "topic_id": {"type": "number", "description": "Discussion assignment ID or discussion topic ID (for graded discussions, use the assignment ID)"},
                    "user_id": {"type": "number", "description": "Student user ID to check"},
                },
                "required": ["topic_id", "user_id"],
            },
        ),

        # File grading tools
        Tool(
            name="grade_file_submission",
            description=(
                "Download and AI-grade a single student's file submission (Word .docx, Excel .xlsx, or PowerPoint .pptx). "
                "Fetches the assignment instructions and rubric from Canvas, extracts the document content, and uses "
                "Claude to evaluate the work against the rubric. Optionally posts the grade back to Canvas."
            ),
            inputSchema={
                "type": "object",
                "properties": {
                    "assignment_id": {"type": "number", "description": "Assignment ID"},
                    "user_id": {"type": "number", "description": "Student user ID"},
                    "auto_update_grade": {
                        "type": "boolean",
                        "description": "If true, automatically post the AI-recommended grade and feedback to Canvas",
                        "default": False,
                    },
                },
                "required": ["assignment_id", "user_id"],
            },
        ),
        Tool(
            name="grade_all_file_submissions",
            description=(
                "Download and AI-grade every student's file submission for an assignment (Word, Excel, or PowerPoint). "
                "For each student who submitted a file, extracts content and grades it against the assignment rubric using Claude. "
                "Returns a full report. Optionally posts grades back to Canvas. Note: may take several minutes for large classes."
            ),
            inputSchema={
                "type": "object",
                "properties": {
                    "assignment_id": {"type": "number", "description": "Assignment ID"},
                    "auto_update_grades": {
                        "type": "boolean",
                        "description": "If true, automatically post AI-recommended grades and feedback to Canvas for all students",
                        "default": False,
                    },
                },
                "required": ["assignment_id"],
            },
        ),
    ]


@app.call_tool()
async def call_tool(name: str, arguments: Any) -> list[TextContent]:
    """Handle tool calls for Canvas operations."""
    try:
        course = get_course()

        # Assignment tools
        if name == "list_assignments":
            assignments = course.get_assignments()
            result = []
            for assignment in assignments:
                result.append(
                    f"ID: {assignment.id}\n"
                    f"Name: {assignment.name}\n"
                    f"Points: {getattr(assignment, 'points_possible', 'N/A')}\n"
                    f"Due: {getattr(assignment, 'due_at', 'No due date')}\n"
                    f"Published: {getattr(assignment, 'published', False)}\n"
                )
            return [TextContent(type="text", text="\n---\n".join(result) if result else "No assignments found")]

        elif name == "create_assignment":
            assignment_data = {
                "name": arguments["name"],
                "points_possible": arguments["points"],
                "published": arguments.get("published", False),
            }
            if "description" in arguments:
                assignment_data["description"] = arguments["description"]
            if "due_at" in arguments:
                assignment_data["due_at"] = arguments["due_at"]

            assignment = course.create_assignment({"assignment": assignment_data})
            return [TextContent(
                type="text",
                text=f"Assignment created successfully!\nID: {assignment.id}\nName: {assignment.name}\nPoints: {assignment.points_possible}"
            )]

        elif name == "get_assignment_submissions":
            assignment = course.get_assignment(arguments["assignment_id"])
            submissions = assignment.get_submissions()
            result = []
            for submission in submissions:
                user = course.get_user(submission.user_id)
                result.append(
                    f"Student: {user.name} (ID: {user.id})\n"
                    f"Score: {getattr(submission, 'score', 'Not graded')}\n"
                    f"Submitted: {getattr(submission, 'submitted_at', 'Not submitted')}\n"
                    f"Status: {getattr(submission, 'workflow_state', 'Unknown')}"
                )
            return [TextContent(type="text", text="\n---\n".join(result) if result else "No submissions found")]

        elif name == "update_grade":
            assignment = course.get_assignment(arguments["assignment_id"])
            submission = assignment.get_submission(arguments["user_id"])

            update_data = {"posted_grade": arguments["grade"]}
            if "comment" in arguments:
                update_data["comment"] = {"text_comment": arguments["comment"]}

            submission.edit(submission=update_data)
            return [TextContent(
                type="text",
                text=f"Grade updated successfully for user {arguments['user_id']}: {arguments['grade']}"
            )]

        # Student tools
        elif name == "list_students":
            students = course.get_users(enrollment_type=["student"])
            result = []
            for student in students:
                result.append(
                    f"ID: {student.id}\n"
                    f"Name: {student.name}\n"
                    f"Email: {getattr(student, 'email', 'N/A')}\n"
                    f"SIS ID: {getattr(student, 'sis_user_id', 'N/A')}"
                )
            return [TextContent(type="text", text="\n---\n".join(result) if result else "No students found")]

        elif name == "get_student_grades":
            user_id = arguments["user_id"]
            user = course.get_user(user_id)
            
            # Get enrollment with grades
            enrollments = course.get_enrollments(user_id=[user_id], include=['current_points', 'current_score'])
            enrollment = None
            for enroll in enrollments:
                if enroll.user_id == user_id and enroll.type == 'StudentEnrollment':
                    enrollment = enroll
                    break
            
            result = [f"Student: {user.name}"]
            
            # Add overall grade if available
            if enrollment:
                current_score = getattr(enrollment, 'grades', {}).get('current_score', 'N/A')
                final_score = getattr(enrollment, 'grades', {}).get('final_score', 'N/A')
                result.append(f"Current Score: {current_score}%")
                result.append(f"Final Score: {final_score}%")
                result.append("")
            
            # Get individual assignment grades
            assignments = course.get_assignments()
            result.append("Assignment Grades:")
            for assignment in assignments:
                try:
                    submission = assignment.get_submission(user_id)
                    score = getattr(submission, 'score', None)
                    points_possible = getattr(assignment, 'points_possible', 0)
                    if score is not None:
                        result.append(f"  {assignment.name}: {score}/{points_possible}")
                    else:
                        result.append(f"  {assignment.name}: Not graded")
                except Exception:
                    result.append(f"  {assignment.name}: No submission")

            return [TextContent(type="text", text="\n".join(result))]

        elif name == "get_grade_summary":
            students = course.get_users(enrollment_type=["student"])
            student_list = list(students)

            if not student_list:
                return [TextContent(type="text", text="No students found in the course")]

            # Get enrollments with grades for all students
            user_ids = [s.id for s in student_list]
            enrollments = course.get_enrollments(
                type=["StudentEnrollment"],
                include=["current_points", "total_scores"]
            )

            # Build a map of user_id to enrollment grades
            enrollment_map = {}
            for enrollment in enrollments:
                if enrollment.type == "StudentEnrollment":
                    enrollment_map[enrollment.user_id] = enrollment

            result = ["=" * 60]
            result.append("STUDENT GRADE SUMMARY")
            result.append("=" * 60)
            result.append("")

            # Sort students by name
            student_list.sort(key=lambda s: s.name)

            for student in student_list:
                enrollment = enrollment_map.get(student.id)

                if enrollment and hasattr(enrollment, 'grades'):
                    grades = enrollment.grades
                    current_score = grades.get('current_score', 'N/A')
                    current_grade = grades.get('current_grade', '')
                    final_score = grades.get('final_score', 'N/A')
                    final_grade = grades.get('final_grade', '')

                    # Format the score display
                    score_display = f"{current_score}%" if current_score != 'N/A' else "N/A"
                    if current_grade:
                        score_display += f" ({current_grade})"

                    final_display = f"{final_score}%" if final_score != 'N/A' else "N/A"
                    if final_grade:
                        final_display += f" ({final_grade})"

                    result.append(f"Student: {student.name}")
                    result.append(f"  ID: {student.id}")
                    result.append(f"  Current Score: {score_display}")
                    result.append(f"  Final Score: {final_display}")
                else:
                    result.append(f"Student: {student.name}")
                    result.append(f"  ID: {student.id}")
                    result.append(f"  Current Score: No grades available")

                result.append("")

            result.append("=" * 60)
            result.append(f"Total Students: {len(student_list)}")
            result.append("=" * 60)

            return [TextContent(type="text", text="\n".join(result))]

        # Announcement tools
        elif name == "list_announcements":
            limit = arguments.get("limit", 10)
            announcements = course.get_discussion_topics(only_announcements=True)
            result = []
            count = 0
            for announcement in announcements:
                if count >= limit:
                    break
                result.append(
                    f"ID: {announcement.id}\n"
                    f"Title: {announcement.title}\n"
                    f"Posted: {getattr(announcement, 'posted_at', 'N/A')}\n"
                    f"Message: {getattr(announcement, 'message', 'N/A')[:200]}..."
                )
                count += 1
            return [TextContent(type="text", text="\n---\n".join(result) if result else "No announcements found")]

        elif name == "create_announcement":
            announcement = course.create_discussion_topic(
                title=arguments["title"],
                message=arguments["message"],
                is_announcement=True,
                published=arguments.get("published", True),
            )
            return [TextContent(
                type="text",
                text=f"Announcement created successfully!\nID: {announcement.id}\nTitle: {announcement.title}"
            )]

        # Module tools
        elif name == "list_modules":
            modules = course.get_modules()
            result = []
            for module in modules:
                result.append(
                    f"ID: {module.id}\n"
                    f"Name: {module.name}\n"
                    f"Published: {getattr(module, 'published', False)}\n"
                    f"Items: {getattr(module, 'items_count', 0)}"
                )
            return [TextContent(type="text", text="\n---\n".join(result) if result else "No modules found")]

        elif name == "get_module_items":
            module = course.get_module(arguments["module_id"])
            items = module.get_module_items()
            result = [f"Module: {module.name}\n"]
            for item in items:
                result.append(
                    f"  - {item.title} ({item.type})"
                )
            return [TextContent(type="text", text="\n".join(result))]

        elif name == "create_module":
            module = course.create_module(
                module={"name": arguments["name"], "published": arguments.get("published", False)}
            )
            return [TextContent(
                type="text",
                text=f"Module created successfully!\nID: {module.id}\nName: {module.name}"
            )]

        # Rubric tools
        elif name == "get_rubric":
            assignment_id = arguments["assignment_id"]
            if assignment_id in RUBRICS:
                rubric = RUBRICS[assignment_id]
                result = []
                result.append(f"RUBRIC: {rubric['assignment_name']}")
                result.append(f"Total Points: {rubric['total_points']}")
                result.append("=" * 60)

                for dimension in rubric["dimensions"]:
                    result.append(f"\n{dimension['name']} (Max: {dimension['max_points']} points)")
                    result.append("-" * 40)
                    # Sort levels by points descending
                    for points in sorted(dimension["level_exemplars"].keys(), key=int, reverse=True):
                        descriptions = dimension["level_exemplars"][points]
                        result.append(f"  {points} pts: {descriptions[0]}")

                result.append("\n" + "=" * 60)
                return [TextContent(type="text", text="\n".join(result))]
            else:
                return [TextContent(type="text", text=f"No rubric found for assignment ID {assignment_id}")]

        # Discussion tools
        elif name == "get_discussion_posts":
            topic_id = arguments["topic_id"]
            user_id_filter = arguments.get("user_id")

            # Try to get the discussion topic - first try as assignment ID, then as topic ID
            topic = None
            try:
                # Try as assignment ID first (for graded discussions)
                assignment = course.get_assignment(topic_id)
                if hasattr(assignment, 'discussion_topic'):
                    actual_topic_id = assignment.discussion_topic['id']
                    topic = course.get_discussion_topic(actual_topic_id)
            except:
                pass

            if not topic:
                # Try as direct discussion topic ID
                try:
                    topic = course.get_discussion_topic(topic_id)
                except:
                    return [TextContent(type="text", text=f"Discussion topic with ID {topic_id} not found. Please verify the ID is correct.")]

            result = []
            result.append(f"DISCUSSION: {topic.title}")
            result.append("=" * 60)

            # Get all entries (posts) for this topic
            entries = topic.get_topic_entries()

            # Build a map of user_id to user name for efficiency
            user_cache = {}

            def get_user_name(uid):
                if uid not in user_cache:
                    try:
                        user = course.get_user(uid)
                        user_cache[uid] = user.name
                    except Exception:
                        user_cache[uid] = f"User {uid}"
                return user_cache[uid]

            def strip_html(text):
                """Remove HTML tags from text."""
                if not text:
                    return ""
                clean = re.sub(r'<[^>]+>', '', text)
                return clean.strip()

            def count_words(text):
                """Count words in text."""
                if not text:
                    return 0
                clean_text = strip_html(text)
                return len(clean_text.split())

            # Process each top-level entry
            for entry in entries:
                entry_user_id = getattr(entry, 'user_id', None)

                # Skip if filtering by user and this isn't their post
                if user_id_filter and entry_user_id != user_id_filter:
                    # But still check replies for this user
                    pass

                user_name = get_user_name(entry_user_id) if entry_user_id else "Unknown"
                message = getattr(entry, 'message', '')
                word_count = count_words(message)
                created_at = getattr(entry, 'created_at', 'Unknown date')

                # Only show this entry if no filter or matches filter
                if not user_id_filter or entry_user_id == user_id_filter:
                    result.append(f"\n{'─' * 60}")
                    result.append(f"STUDENT: {user_name} (ID: {entry_user_id})")
                    result.append(f"Posted: {created_at}")
                    result.append(f"Word Count: {word_count}")
                    result.append(f"{'─' * 40}")
                    result.append(f"INITIAL POST:")
                    result.append(strip_html(message) if message else "(No content)")

                # Get replies to this entry
                try:
                    replies = entry.get_replies()
                    reply_count = 0
                    for reply in replies:
                        reply_user_id = getattr(reply, 'user_id', None)

                        # If filtering, only show replies from that user
                        if user_id_filter and reply_user_id != user_id_filter:
                            continue

                        reply_user_name = get_user_name(reply_user_id) if reply_user_id else "Unknown"
                        reply_message = getattr(reply, 'message', '')
                        reply_word_count = count_words(reply_message)
                        reply_created_at = getattr(reply, 'created_at', 'Unknown date')

                        if not user_id_filter or entry_user_id == user_id_filter:
                            # Show as reply under the original post
                            result.append(f"\n  ↳ REPLY from {reply_user_name} (ID: {reply_user_id})")
                            result.append(f"    Posted: {reply_created_at} | Words: {reply_word_count}")
                            result.append(f"    {strip_html(reply_message)[:500]}{'...' if len(strip_html(reply_message)) > 500 else ''}")
                        elif reply_user_id == user_id_filter:
                            # This is a reply FROM the filtered user to someone else's post
                            result.append(f"\n{'─' * 60}")
                            result.append(f"REPLY by {reply_user_name} (ID: {reply_user_id})")
                            result.append(f"In response to: {user_name}'s post")
                            result.append(f"Posted: {reply_created_at} | Words: {reply_word_count}")
                            result.append(f"{'─' * 40}")
                            result.append(strip_html(reply_message) if reply_message else "(No content)")

                        reply_count += 1

                    if not user_id_filter or entry_user_id == user_id_filter:
                        if reply_count == 0:
                            result.append(f"\n  (No replies to this post)")
                except Exception as e:
                    logger.debug(f"Could not get replies for entry: {e}")

            result.append(f"\n{'=' * 60}")
            return [TextContent(type="text", text="\n".join(result))]

        # AI Detection tools
        elif name == "check_ai_writing":
            topic_id = arguments["topic_id"]
            user_id = arguments["user_id"]

            # Try to get the discussion topic - first try as assignment ID, then as topic ID
            topic = None
            try:
                # Try as assignment ID first (for graded discussions)
                assignment = course.get_assignment(topic_id)
                if hasattr(assignment, 'discussion_topic'):
                    actual_topic_id = assignment.discussion_topic['id']
                    topic = course.get_discussion_topic(actual_topic_id)
            except:
                pass

            if not topic:
                # Try as direct discussion topic ID
                try:
                    topic = course.get_discussion_topic(topic_id)
                except:
                    return [TextContent(type="text", text=f"Discussion topic with ID {topic_id} not found. Please verify the ID is correct.")]

            entries = topic.get_topic_entries()

            def strip_html(text):
                """Remove HTML tags from text."""
                if not text:
                    return ""
                clean = re.sub(r'<[^>]+>', '', text)
                return clean.strip()

            # Find the student's post
            student_post = None
            student_name = None
            for entry in entries:
                if getattr(entry, 'user_id', None) == user_id:
                    student_post = strip_html(getattr(entry, 'message', ''))
                    try:
                        user = course.get_user(user_id)
                        student_name = user.name
                    except Exception:
                        student_name = f"User {user_id}"
                    break

            if not student_post:
                return [TextContent(type="text", text=f"No post found for user ID {user_id} in this discussion.")]

            # Use Claude to analyze the text
            api_key = os.getenv("ANTHROPIC_API_KEY")
            if not api_key:
                return [TextContent(type="text", text="Error: ANTHROPIC_API_KEY environment variable not set.")]

            client = anthropic.Anthropic(api_key=api_key)

            analysis_prompt = f"""Analyze the following student discussion post and assess whether it was likely written by AI (like ChatGPT, Claude, etc.) or by a human student.

Consider these factors:
1. Writing style - Is it overly formal, generic, or lacks personal voice?
2. Structure - Does it follow a rigid template (intro, body, conclusion) typical of AI?
3. Phrases - Are there AI-typical phrases like "In conclusion," "It's important to note," "Overall," "That being said"?
4. Personal details - Does it include specific personal anecdotes or experiences that feel authentic?
5. Errors - Does it have natural human errors or typos that AI typically wouldn't make?
6. Vocabulary - Is the vocabulary consistent with a student or unusually sophisticated?
7. Hedging language - Excessive use of "I believe," "In my opinion," "I think" can be AI indicators

STUDENT POST:
\"\"\"
{student_post}
\"\"\"

Provide your analysis in this format:
LIKELIHOOD: [Low/Medium/High] likelihood of AI-generated content
CONFIDENCE: [Low/Medium/High] confidence in this assessment

KEY INDICATORS:
- [List 3-5 specific observations from the text]

SUMMARY: [2-3 sentence summary of your analysis]"""

            message = client.messages.create(
                model="claude-sonnet-4-20250514",
                max_tokens=1024,
                messages=[
                    {"role": "user", "content": analysis_prompt}
                ]
            )

            analysis_result = message.content[0].text

            result = []
            result.append(f"AI WRITING ANALYSIS")
            result.append("=" * 60)
            result.append(f"Student: {student_name} (ID: {user_id})")
            result.append(f"Discussion: {topic.title}")
            result.append(f"Post Length: {len(student_post.split())} words")
            result.append("=" * 60)
            result.append("")
            result.append(analysis_result)
            result.append("")
            result.append("=" * 60)
            result.append("NOTE: This is an AI-based assessment and should be used")
            result.append("as one factor among many when evaluating student work.")
            result.append("=" * 60)

            return [TextContent(type="text", text="\n".join(result))]

        # File grading tools
        elif name == "grade_file_submission":
            assignment_id = arguments["assignment_id"]
            user_id = arguments["user_id"]
            auto_update = arguments.get("auto_update_grade", False)

            assignment = course.get_assignment(assignment_id)
            submission = assignment.get_submission(user_id)
            user = course.get_user(user_id)

            rubric_text = _rubric_text_from_dict(assignment_id) or _get_canvas_rubric(assignment)
            assignment_description = _strip_html(getattr(assignment, "description", "") or "")
            points_possible = float(getattr(assignment, "points_possible", 100) or 100)

            result = [
                f"GRADING: {user.name}",
                f"Assignment: {assignment.name}",
                f"Points Possible: {points_possible}",
                "=" * 60,
            ]

            grade_lines = await _grade_submission(
                course=course,
                assignment=assignment,
                submission=submission,
                student_name=user.name,
                rubric_text=rubric_text,
                assignment_description=assignment_description,
                points_possible=points_possible,
                auto_update=auto_update,
            )
            result.extend(grade_lines)
            return [TextContent(type="text", text="\n".join(result))]

        elif name == "grade_all_file_submissions":
            assignment_id = arguments["assignment_id"]
            auto_update = arguments.get("auto_update_grades", False)

            assignment = course.get_assignment(assignment_id)
            submissions = list(assignment.get_submissions())

            rubric_text = _rubric_text_from_dict(assignment_id) or _get_canvas_rubric(assignment)
            assignment_description = _strip_html(getattr(assignment, "description", "") or "")
            points_possible = float(getattr(assignment, "points_possible", 100) or 100)

            result = [
                "GRADING ALL SUBMISSIONS",
                f"Assignment: {assignment.name}",
                f"Points Possible: {points_possible}",
                f"Total Submissions Found: {len(submissions)}",
                "=" * 60,
            ]

            graded_count = 0
            skipped_count = 0

            for submission in submissions:
                attachments = getattr(submission, "attachments", []) or []
                if not attachments:
                    skipped_count += 1
                    continue

                try:
                    user = course.get_user(submission.user_id)
                    student_name = user.name
                except Exception:
                    student_name = f"User {submission.user_id}"

                result.append(f"\nStudent: {student_name} (ID: {submission.user_id})")
                result.append("-" * 60)

                grade_lines = await _grade_submission(
                    course=course,
                    assignment=assignment,
                    submission=submission,
                    student_name=student_name,
                    rubric_text=rubric_text,
                    assignment_description=assignment_description,
                    points_possible=points_possible,
                    auto_update=auto_update,
                )
                result.extend(grade_lines)
                result.append("=" * 60)
                graded_count += 1

            result.append("")
            result.append(f"SUMMARY: Graded {graded_count} student(s), skipped {skipped_count} (no file submission).")
            if auto_update:
                result.append("Grades have been posted to Canvas.")
            return [TextContent(type="text", text="\n".join(result))]

        else:
            return [TextContent(type="text", text=f"Unknown tool: {name}")]

    except CanvasException as e:
        logger.error(f"Canvas API error: {e}")
        return [TextContent(type="text", text=f"Canvas API error: {str(e)}")]
    except Exception as e:
        logger.error(f"Error executing tool {name}: {e}")
        return [TextContent(type="text", text=f"Error: {str(e)}")]


async def main():
    """Run the MCP server."""
    from mcp.server.stdio import stdio_server

    async with stdio_server() as (read_stream, write_stream):
        await app.run(
            read_stream,
            write_stream,
            app.create_initialization_options(),
        )


if __name__ == "__main__":
    import asyncio
    asyncio.run(main())
